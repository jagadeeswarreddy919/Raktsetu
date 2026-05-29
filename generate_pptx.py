import os
import sys
import subprocess

# Ensure python-pptx is installed
try:
    import pptx
except ImportError:
    print("python-pptx is missing. Installing it now...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Use 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # Palette definition (Slate & Crimson Theme)
    BG_COLOR = RGBColor(15, 23, 42)        # Slate 900
    CARD_BG = RGBColor(30, 41, 59)         # Slate 800
    BORDER_COLOR = RGBColor(51, 65, 85)     # Slate 700
    TEXT_MAIN = RGBColor(255, 255, 255)     # White
    TEXT_MUTED = RGBColor(148, 163, 184)   # Slate 400
    CRIMSON = RGBColor(239, 68, 68)        # Crimson Red Accent
    EMERALD = RGBColor(16, 185, 129)       # Emerald Green Accent
    AMBER = RGBColor(245, 158, 11)         # Amber Accent
    
    def apply_slide_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
    def add_header(slide, title_text, category_text="RAKTSETU PLATFORM"):
        # Category tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.3))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        cat_tf.margin_left = cat_tf.margin_top = cat_tf.margin_right = cat_tf.margin_bottom = 0
        p_cat = cat_tf.paragraphs[0]
        run_cat = p_cat.add_run()
        run_cat.text = category_text.upper()
        run_cat.font.name = "Arial"
        run_cat.font.size = Pt(10)
        run_cat.font.bold = True
        run_cat.font.color.rgb = CRIMSON
        
        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.733), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.margin_left = title_tf.margin_top = title_tf.margin_right = title_tf.margin_bottom = 0
        p_title = title_tf.paragraphs[0]
        run_title = p_title.add_run()
        run_title.text = title_text
        run_title.font.name = "Arial"
        run_title.font.size = Pt(28)
        run_title.font.bold = True
        run_title.font.color.rgb = TEXT_MAIN
        
    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Sleek, Premium Cover)
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    apply_slide_bg(slide1)
    
    # Large Decorative Left Highlight
    left_accent = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    left_accent.fill.solid()
    left_accent.fill.fore_color.rgb = CRIMSON
    left_accent.line.fill.background()
    
    # Main Title & Subtitle Frame
    main_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.0), Inches(4.0))
    tf1 = main_box.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_top = tf1.margin_right = tf1.margin_bottom = 0
    
    p_main = tf1.paragraphs[0]
    r_main = p_main.add_run()
    r_main.text = "RAKTSETU"
    r_main.font.name = "Arial"
    r_main.font.size = Pt(64)
    r_main.font.bold = True
    r_main.font.color.rgb = TEXT_MAIN
    
    p_sub = tf1.add_paragraph()
    p_sub.space_before = Pt(12)
    r_sub = p_sub.add_run()
    r_sub.text = "Next-Generation Smart Blood Management & P2P Bridging Platform"
    r_sub.font.name = "Arial"
    r_sub.font.size = Pt(20)
    r_sub.font.bold = True
    r_sub.font.color.rgb = CRIMSON
    
    p_desc = tf1.add_paragraph()
    p_desc.space_before = Pt(24)
    r_desc = p_desc.add_run()
    r_desc.text = "An intelligent, high-throughput microservices architecture connecting donors, recipients, and hospitals with real-time communications and machine learning analytics."
    r_desc.font.name = "Arial"
    r_desc.font.size = Pt(14)
    r_desc.font.color.rgb = TEXT_MUTED
    
    # Bottom Footer
    footer_box = slide1.shapes.add_textbox(Inches(1.2), Inches(6.2), Inches(11.0), Inches(0.5))
    tf_f = footer_box.text_frame
    tf_f.word_wrap = True
    p_f = tf_f.paragraphs[0]
    r_f = p_f.add_run()
    r_f.text = "System Architecture & Comprehensive Overview  |  Technical Slide Deck"
    r_f.font.name = "Arial"
    r_f.font.size = Pt(11)
    r_f.font.color.rgb = TEXT_MUTED
    
    # ----------------------------------------------------
    # SLIDE 2: The Core Challenges (Problem Statement)
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    apply_slide_bg(slide2)
    add_header(slide2, "The Critical Challenges in Blood Donation & Matching", "THE PROBLEM STATEMENT")
    
    # 3 Column Layout for Problems
    probs = [
        {
            "title": "Information & Supply Silos",
            "desc": "Blood banks, hospitals, and active donors operate in highly fragmented silos. Real-time availability of specific blood groups is extremely restricted, resulting in critical communication delays during medical emergencies.",
            "color": CRIMSON,
            "badge": "SILOS"
        },
        {
            "title": "Inefficient Donor Matching",
            "desc": "Emergency requests rely heavily on manual calls and social media broadcasting. Locating verified, eligible, and physically nearby donors in real-time is a chaotic process lacking scientific ranking or proximity filtering.",
            "color": AMBER,
            "badge": "FRICTION"
        },
        {
            "title": "Seasonal Shortage Vulnerability",
            "desc": "Hospitals struggle with predictive inventory management. Without forecasting regional blood demand spikes and donation dips, administrations fail to take proactive measures before deficit situations arise.",
            "color": CRIMSON,
            "badge": "DEFICITS"
        }
    ]
    
    card_width = Inches(3.64)
    card_gap = Inches(0.4)
    card_y = Inches(2.0)
    card_h = Inches(4.5)
    
    for i, prob in enumerate(probs):
        card_x = Inches(0.8) + i * (card_width + card_gap)
        
        # Card Background
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, card_width, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1.5)
        
        # Card Header Indicator Bar
        bar = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, card_x, card_y, card_width, Inches(0.12))
        bar.fill.solid()
        bar.fill.fore_color.rgb = prob["color"]
        bar.line.fill.background()
        
        # Card Text Box
        tbox = slide2.shapes.add_textbox(card_x + Inches(0.25), card_y + Inches(0.3), card_width - Inches(0.5), card_h - Inches(0.5))
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        # Badge
        pb = tf.paragraphs[0]
        rb = pb.add_run()
        rb.text = f"[{prob['badge']}]"
        rb.font.name = "Arial"
        rb.font.size = Pt(10)
        rb.font.bold = True
        rb.font.color.rgb = prob["color"]
        
        # Title
        pt = tf.add_paragraph()
        pt.space_before = Pt(12)
        rt = pt.add_run()
        rt.text = prob["title"]
        rt.font.name = "Arial"
        rt.font.size = Pt(20)
        rt.font.bold = True
        rt.font.color.rgb = TEXT_MAIN
        
        # Desc
        pd = tf.add_paragraph()
        pd.space_before = Pt(16)
        rd = pd.add_run()
        rd.text = prob["desc"]
        rd.font.name = "Arial"
        rd.font.size = Pt(13)
        rd.font.color.rgb = TEXT_MUTED
        
    # ----------------------------------------------------
    # SLIDE 3: The RaktSetu Solution (Core Value Props)
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    apply_slide_bg(slide3)
    add_header(slide3, "RaktSetu: The Intelligent Blood Bridge", "THE CORE SOLUTION")
    
    # 4 Column Grid of Solutions
    sols = [
        {
            "title": "Instant P2P Chat",
            "desc": "Seamlessly connects donor and recipient directly in a secure, real-time chat workspace. Facilitates immediate coordination without publicizing sensitive contact details.",
            "icon_color": CRIMSON,
            "badge": "REAL-TIME"
        },
        {
            "title": "AI Match & Rank",
            "desc": "ML algorithms parse geographical factors, active availability, verification status, and historical eligibility to rank and list the ideal matching donors within seconds.",
            "icon_color": EMERALD,
            "badge": "SMART FILTER"
        },
        {
            "title": "Predictive Shortage",
            "desc": "Empowers hospitals with Scikit-Learn based intelligence forecasting blood inventory deficit levels dynamically based on local historical data and seasonal trends.",
            "icon_color": AMBER,
            "badge": "FORECASTING"
        },
        {
            "title": "Gamified Rewards",
            "desc": "Encourages a continuous culture of blood donations through a point-based gamified structure, milestone achievement tracker, and verified donor validation trust badges.",
            "icon_color": EMERALD,
            "badge": "RETENTION"
        }
    ]
    
    sol_width = Inches(2.68)
    sol_gap = Inches(0.33)
    
    for i, sol in enumerate(sols):
        sol_x = Inches(0.8) + i * (sol_width + sol_gap)
        
        # Solution Card
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sol_x, Inches(2.0), sol_width, Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1.5)
        
        # Indicator Dot / Accent
        dot = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, sol_x, Inches(2.0), sol_width, Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = sol["icon_color"]
        dot.line.fill.background()
        
        tbox = slide3.shapes.add_textbox(sol_x + Inches(0.2), Inches(2.3), sol_width - Inches(0.4), Inches(4.0))
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        # Category Tag
        pb = tf.paragraphs[0]
        rb = pb.add_run()
        rb.text = sol["badge"]
        rb.font.name = "Arial"
        rb.font.size = Pt(9)
        rb.font.bold = True
        rb.font.color.rgb = sol["icon_color"]
        
        # Title
        pt = tf.add_paragraph()
        pt.space_before = Pt(8)
        rt = pt.add_run()
        rt.text = sol["title"]
        rt.font.name = "Arial"
        rt.font.size = Pt(18)
        rt.font.bold = True
        rt.font.color.rgb = TEXT_MAIN
        
        # Desc
        pd = tf.add_paragraph()
        pd.space_before = Pt(12)
        rd = pd.add_run()
        rd.text = sol["desc"]
        rd.font.name = "Arial"
        rd.font.size = Pt(12)
        rd.font.color.rgb = TEXT_MUTED
        
    # ----------------------------------------------------
    # SLIDE 4: High-Level System Architecture
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    apply_slide_bg(slide4)
    add_header(slide4, "Robust High-Performance Microservices Architecture", "SYSTEM ARCHITECTURE")
    
    # Draw Architecture layers
    layers = [
        {"name": "1. PREMIUM CLIENT PORTAL", "tech": "React 18 / Redux Toolkit / CSS Glassmorphism", "desc": "Role-specific intuitive portals for Admin, Donor, Recipient, and Hospital. Dynamic dashboard states & localized routing.", "color": CRIMSON},
        {"name": "2. API GATEWAY & CORE BACKEND", "tech": "Node.js / Express.js / MongoDB Database", "desc": "Handles core API requests, authenticates sessions with secure JWTs, coordinates chat sockets, and manages persistent user schema structures.", "color": AMBER},
        {"name": "3. PREDICTIVE ML ENGINE", "tech": "FastAPI Microservice / Python / Scikit-Learn", "desc": "Lightweight, highly optimized microservice supplying regional shortage predictions & applying multidimensional ranking of eligible donors.", "color": EMERALD},
        {"name": "4. HIGH-THROUGHPUT ANALYTICS", "tech": "Java Spring Boot / Scheduled Tasks Consolidation", "desc": "Consolidates large-volume system logs asynchronously, monitors queue parameters, and evaluates metrics under extreme latency targets.", "color": TEXT_MAIN}
    ]
    
    for i, layer in enumerate(layers):
        layer_y = Inches(1.95) + i * Inches(1.25)
        
        # Layer border box
        box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), layer_y, Inches(11.733), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = BORDER_COLOR
        box.line.width = Pt(1.2)
        
        # Left color strip
        strip = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), layer_y, Inches(0.15), Inches(1.1))
        strip.fill.solid()
        strip.fill.fore_color.rgb = layer["color"]
        strip.line.fill.background()
        
        # Content
        tbox = slide4.shapes.add_textbox(Inches(1.2), layer_y + Inches(0.12), Inches(11.0), Inches(0.9))
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        # Title and tech stack
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = layer["name"]
        r1.font.name = "Arial"
        r1.font.size = Pt(14)
        r1.font.bold = True
        r1.font.color.rgb = layer["color"]
        
        r1_tech = p1.add_run()
        r1_tech.text = f"   |   {layer['tech']}"
        r1_tech.font.name = "Arial"
        r1_tech.font.size = Pt(12)
        r1_tech.font.bold = True
        r1_tech.font.color.rgb = TEXT_MAIN
        
        # Desc
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        r2 = p2.add_run()
        r2.text = layer["desc"]
        r2.font.name = "Arial"
        r2.font.size = Pt(11.5)
        r2.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 5: Role-Based Dashboards & Workflows
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    apply_slide_bg(slide5)
    add_header(slide5, "Targeted Operations: Multi-Role Dashboard States", "USER DASHBOARD SYSTEM")
    
    roles = [
        {
            "name": "Donor Dashboard",
            "points": [
                "Gamified Reward Points & verified donor validation badges.",
                "Real-time toggles for donor active availability status.",
                "Comprehensive medical documentation & upload interface.",
                "Upcoming campaign registrations and historic timeline."
            ],
            "accent": EMERALD
        },
        {
            "name": "Hospital Dashboard",
            "points": [
                "Real-time blood stock levels and storage management.",
                "Predictive analytics panel predicting regional shortage levels.",
                "Advanced donor directory with proximity filtering.",
                "Official campaign builder and notification broadcasts."
            ],
            "accent": CRIMSON
        },
        {
            "name": "Recipient & Admin",
            "points": [
                "Fast, clean blood emergency requests generation wizard.",
                "P2P direct chat hub for immediate donor coordinate sessions.",
                "Admin audit logs, campaign moderation, and user checks.",
                "Comprehensive visual metrics representing platform status."
            ],
            "accent": AMBER
        }
    ]
    
    role_width = Inches(3.64)
    role_gap = Inches(0.4)
    
    for i, role in enumerate(roles):
        role_x = Inches(0.8) + i * (role_width + role_gap)
        
        # Role Card
        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, role_x, Inches(2.0), role_width, Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1.5)
        
        # Color line
        line = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, role_x, Inches(2.0), role_width, Inches(0.12))
        line.fill.solid()
        line.fill.fore_color.rgb = role["accent"]
        line.line.fill.background()
        
        tbox = slide5.shapes.add_textbox(role_x + Inches(0.25), Inches(2.3), role_width - Inches(0.5), Inches(4.0))
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        # Title
        p_title = tf.paragraphs[0]
        r_title = p_title.add_run()
        r_title.text = role["name"]
        r_title.font.name = "Arial"
        r_title.font.size = Pt(20)
        r_title.font.bold = True
        r_title.font.color.rgb = TEXT_MAIN
        
        # List of Points
        for pt in role["points"]:
            p_pt = tf.add_paragraph()
            p_pt.space_before = Pt(12)
            p_pt.level = 0
            
            # Bullet point accent
            bullet = p_pt.add_run()
            bullet.text = "• "
            bullet.font.bold = True
            bullet.font.color.rgb = role["accent"]
            
            r_pt = p_pt.add_run()
            r_pt.text = pt
            r_pt.font.name = "Arial"
            r_pt.font.size = Pt(11.5)
            r_pt.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 6: AI Engine & Machine Learning Microservice
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    apply_slide_bg(slide6)
    add_header(slide6, "Deep Dive: Smart Donor Ranking & Shortage Prediction", "PREDICTIVE AI & FASTAPI")
    
    # 2 Column layout
    # Left Column: Shortage Prediction
    left_x = Inches(0.8)
    col_width = Inches(5.66)
    
    card_l = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, Inches(2.0), col_width, Inches(4.5))
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = CARD_BG
    card_l.line.color.rgb = BORDER_COLOR
    card_l.line.width = Pt(1.5)
    
    bar_l = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, Inches(2.0), col_width, Inches(0.12))
    bar_l.fill.solid()
    bar_l.fill.fore_color.rgb = CRIMSON
    bar_l.line.fill.background()
    
    tbox_l = slide6.shapes.add_textbox(left_x + Inches(0.3), Inches(2.3), col_width - Inches(0.6), Inches(4.0))
    tf_l = tbox_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_top = tf_l.margin_right = tf_l.margin_bottom = 0
    
    p_lt = tf_l.paragraphs[0]
    r_lt = p_lt.add_run()
    r_lt.text = "1. Regional Shortage Prediction"
    r_lt.font.name = "Arial"
    r_lt.font.size = Pt(20)
    r_lt.font.bold = True
    r_lt.font.color.rgb = TEXT_MAIN
    
    features_l = [
        ("API Endpoint", "/api/ai/predict-shortage (POST)"),
        ("Model Inputs", "Month, State Code, Pincode Factor, Current Inventory, & Previous Month Demand."),
        ("ML Algorithm", "Scikit-Learn Regression pipeline evaluating stock density values."),
        ("Output Signal", "Quantifiable shortage index representing high-deficit risk vs safe-buffer margins.")
    ]
    
    for label, desc in features_l:
        p = tf_l.add_paragraph()
        p.space_before = Pt(10)
        
        rl = p.add_run()
        rl.text = f"{label}: "
        rl.font.bold = True
        rl.font.size = Pt(12)
        rl.font.color.rgb = CRIMSON
        
        rd = p.add_run()
        rd.text = desc
        rd.font.size = Pt(12)
        rd.font.color.rgb = TEXT_MUTED

    # Right Column: Donor Ranking
    right_x = Inches(6.86)
    card_r = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, Inches(2.0), col_width, Inches(4.5))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = CARD_BG
    card_r.line.color.rgb = BORDER_COLOR
    card_r.line.width = Pt(1.5)
    
    bar_r = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, Inches(2.0), col_width, Inches(0.12))
    bar_r.fill.solid()
    bar_r.fill.fore_color.rgb = EMERALD
    bar_r.line.fill.background()
    
    tbox_r = slide6.shapes.add_textbox(right_x + Inches(0.3), Inches(2.3), col_width - Inches(0.6), Inches(4.0))
    tf_r = tbox_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_top = tf_r.margin_right = tf_r.margin_bottom = 0
    
    p_rt = tf_r.paragraphs[0]
    r_rt = p_rt.add_run()
    r_rt.text = "2. Multi-Criteria Donor Ranking"
    r_rt.font.name = "Arial"
    r_rt.font.size = Pt(20)
    r_rt.font.bold = True
    r_rt.font.color.rgb = TEXT_MAIN
    
    features_r = [
        ("API Endpoint", "/api/ai/rank-donors (POST)"),
        ("Geographical Match", "Evaluates exact proximity matches across State, District, City, Area, & Pincode."),
        ("Engagement Metrics", "Injects user verified checks, last donation date validation (safeguarding safety gaps), and gamified points."),
        ("Result Output", "A sorted catalog of active available donors ranked on real availability score, preventing unnecessary notification fatigue.")
    ]
    
    for label, desc in features_r:
        p = tf_r.add_paragraph()
        p.space_before = Pt(10)
        
        rl = p.add_run()
        rl.text = f"{label}: "
        rl.font.bold = True
        rl.font.size = Pt(12)
        rl.font.color.rgb = EMERALD
        
        rd = p.add_run()
        rd.text = desc
        rd.font.size = Pt(12)
        rd.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 7: High-Throughput Java Analytics Microservice
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    apply_slide_bg(slide7)
    add_header(slide7, "Asynchronous Log Consolidation & Microservice Monitoring", "SPRING BOOT ANALYTICS ENGINE")
    
    # Central Details Card
    card_c = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(11.733), Inches(4.5))
    card_c.fill.solid()
    card_c.fill.fore_color.rgb = CARD_BG
    card_c.line.color.rgb = BORDER_COLOR
    card_c.line.width = Pt(1.5)
    
    bar_c = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(11.733), Inches(0.12))
    bar_c.fill.solid()
    bar_c.fill.fore_color.rgb = AMBER
    bar_c.line.fill.background()
    
    tbox_c = slide7.shapes.add_textbox(Inches(1.2), Inches(2.3), Inches(11.0), Inches(4.0))
    tf_c = tbox_c.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
    
    p_ct = tf_c.paragraphs[0]
    r_ct = p_ct.add_run()
    r_ct.text = "Microservice Specifications: /api/analytics"
    r_ct.font.name = "Arial"
    r_ct.font.size = Pt(22)
    r_ct.font.bold = True
    r_ct.font.color.rgb = TEXT_MAIN
    
    features_c = [
        ("Background Consolidation", "Configured with Spring Scheduling tasks consolidation, running at fixed 15-second cycles to emulate heavy logging clusters and compress operational metrics."),
        ("Simulated Message Queue", "Monitors simulated queue performance indices (average latency stable at 42.5ms, processed jobs: ~1420+, generated reports: ~385+)."),
        ("System Indicators", "Exposes endpoints for system summary (/system-summary) and monthly donation vs demand distribution ratios (/monthly-distribution) for visualization charts."),
        ("Framework Stack", "Java 17, Spring Boot 3.2.5, Web modules, scheduling, cross-origin allowed REST structures for backend gateway synchronization.")
    ]
    
    for label, desc in features_c:
        p = tf_c.add_paragraph()
        p.space_before = Pt(14)
        
        rl = p.add_run()
        rl.text = f"{label}   —   "
        rl.font.bold = True
        rl.font.size = Pt(13.5)
        rl.font.color.rgb = AMBER
        
        rd = p.add_run()
        rd.text = desc
        rd.font.size = Pt(13)
        rd.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 8: Premium Stack & Strategic Roadmap
    # ----------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    apply_slide_bg(slide8)
    add_header(slide8, "Future Outlook, Value Delivery & Technology Stack Summary", "ROADMAP & STACK SUMMARY")
    
    # 2 Column layout
    # Left: Technology Stack Summary
    left_x = Inches(0.8)
    col_width = Inches(5.66)
    
    card_l8 = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, Inches(2.0), col_width, Inches(4.5))
    card_l8.fill.solid()
    card_l8.fill.fore_color.rgb = CARD_BG
    card_l8.line.color.rgb = BORDER_COLOR
    card_l8.line.width = Pt(1.5)
    
    bar_l8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, Inches(2.0), col_width, Inches(0.12))
    bar_l8.fill.solid()
    bar_l8.fill.fore_color.rgb = EMERALD
    bar_l8.line.fill.background()
    
    tbox_l8 = slide8.shapes.add_textbox(left_x + Inches(0.3), Inches(2.3), col_width - Inches(0.6), Inches(4.0))
    tf_l8 = tbox_l8.text_frame
    tf_l8.word_wrap = True
    tf_l8.margin_left = tf_l8.margin_top = tf_l8.margin_right = tf_l8.margin_bottom = 0
    
    p_lt8 = tf_l8.paragraphs[0]
    r_lt8 = p_lt8.add_run()
    r_lt8.text = "Complete Technology Stack"
    r_lt8.font.name = "Arial"
    r_lt8.font.size = Pt(20)
    r_lt8.font.bold = True
    r_lt8.font.color.rgb = TEXT_MAIN
    
    stacks = [
        ("Frontend client", "React, Redux Toolkit, React-Router-Dom v6, Custom Glassmorphism Theme"),
        ("Gateway & API", "Node.js, Express, MongoDB (Mongoose schemas), Socket.io (real-time chat)"),
        ("AI Microservice", "FastAPI (Python), Scikit-Learn, Pydantic, Uvicorn Server"),
        ("Analytics", "Java 17, Spring Boot 3.2.5, Spring Scheduler, Maven dependencies"),
        ("Process Control", "PM2 Ecosystem configuration, Docker-Compose, Powershell Launchers")
    ]
    
    for label, desc in stacks:
        p = tf_l8.add_paragraph()
        p.space_before = Pt(8)
        
        rl = p.add_run()
        rl.text = f"{label}: "
        rl.font.bold = True
        rl.font.size = Pt(11.5)
        rl.font.color.rgb = EMERALD
        
        rd = p.add_run()
        rd.text = desc
        rd.font.size = Pt(11.5)
        rd.font.color.rgb = TEXT_MUTED

    # Right: Roadmap & Future Scale
    right_x = Inches(6.86)
    card_r8 = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, Inches(2.0), col_width, Inches(4.5))
    card_r8.fill.solid()
    card_r8.fill.fore_color.rgb = CARD_BG
    card_r8.line.color.rgb = BORDER_COLOR
    card_r8.line.width = Pt(1.5)
    
    bar_r8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, Inches(2.0), col_width, Inches(0.12))
    bar_r8.fill.solid()
    bar_r8.fill.fore_color.rgb = CRIMSON
    bar_r8.line.fill.background()
    
    tbox_r8 = slide8.shapes.add_textbox(right_x + Inches(0.3), Inches(2.3), col_width - Inches(0.6), Inches(4.0))
    tf_r8 = tbox_r8.text_frame
    tf_r8.word_wrap = True
    tf_r8.margin_left = tf_r8.margin_top = tf_r8.margin_right = tf_r8.margin_bottom = 0
    
    p_rt8 = tf_r8.paragraphs[0]
    r_rt8 = p_rt8.add_run()
    r_rt8.text = "Strategic Roadmap & Scaling"
    r_rt8.font.name = "Arial"
    r_rt8.font.size = Pt(20)
    r_rt8.font.bold = True
    r_rt8.font.color.rgb = TEXT_MAIN
    
    roadmaps = [
        ("Local API Integration", "Sync directly with regional health authorities and private hospital inventory management software arrays."),
        ("Automated Alerts", "Integrate automated cellular SMS/WhatsApp gateways for critical localized blood type shortage announcements."),
        ("Advanced ML Models", "Upgrade current Scikit-Learn pipelines to Deep Learning LSTM networks for long-term seasonal forecasting."),
        ("Native Mobile App", "Build native iOS and Android application platforms to implement geofenced emergency notification alerts.")
    ]
    
    for label, desc in roadmaps:
        p = tf_r8.add_paragraph()
        p.space_before = Pt(8)
        
        rl = p.add_run()
        rl.text = f"{label}: "
        rl.font.bold = True
        rl.font.size = Pt(11.5)
        rl.font.color.rgb = CRIMSON
        
        rd = p.add_run()
        rd.text = desc
        rd.font.size = Pt(11.5)
        rd.font.color.rgb = TEXT_MUTED

    # Save to file
    output_filename = "RaktSetu_Presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as '{output_filename}'!")

if __name__ == "__main__":
    create_presentation()
